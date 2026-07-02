from __future__ import annotations

from io import BytesIO

import pytest

from core.export.base import ExportPayload
from core.export.factory import ExporterFactory
from core.registry import UnknownPluginError


@pytest.mark.asyncio
async def test_markdown_conversation_export() -> None:
    exporter = ExporterFactory.create("markdown")
    result = await exporter.export(
        payload=ExportPayload(
            title="Sales Q&A",
            content_type="conversation",
            data={"messages": [{"role": "user", "content": "revenue?"}, {"role": "assistant", "content": "$1M"}]},
        )
    )
    assert result.filename.endswith(".md")
    assert b"revenue" in result.data


@pytest.mark.asyncio
async def test_pdf_dashboard_export() -> None:
    exporter = ExporterFactory.create("pdf")
    result = await exporter.export(
        payload=ExportPayload(
            title="Ops Dashboard",
            content_type="dashboard",
            data={
                "cards": [
                    {
                        "title": "Customers",
                        "response": {
                            "response_type": "data_table",
                            "data": {
                                "columns": ["id", "name"],
                                "rows": [[1, "Alice"], [2, "Bob"]],
                            },
                        },
                    }
                ]
            },
        )
    )
    assert result.filename.endswith(".pdf")
    assert result.data.startswith(b"%PDF")
    assert len(result.data) > 500


@pytest.mark.asyncio
async def test_format_response_text_data_table() -> None:
    from core.export.response_render import format_response_text

    text = format_response_text(
        {
            "response_type": "data_table",
            "data": {"columns": ["region", "total"], "rows": [["APAC", 100]]},
        }
    )
    assert "region" in text
    assert "APAC" in text
    assert "data_table" not in text


@pytest.mark.asyncio
async def test_pptx_exporter_registered() -> None:
    exporter = ExporterFactory.create("pptx")
    result = await exporter.export(
        payload=ExportPayload(
            title="Deck",
            content_type="dashboard",
            data={"cards": [{"title": "Slide 1", "summary": "hello"}]},
        )
    )
    assert result.filename.endswith(".pptx")
    assert len(result.data) > 100


def test_exporter_registry_keys() -> None:
    keys = ExporterFactory.keys()
    assert "markdown" in keys
    assert "pdf" in keys
    assert "pptx" in keys
    with pytest.raises(UnknownPluginError):
        ExporterFactory.create("docx")


def test_sanitize_filename_strips_unsafe_chars() -> None:
    from core.export.response_render import sanitize_filename

    assert sanitize_filename('Q3 "Revenue" / Report?.txt') == "Q3_Revenue_Report_.txt"
    assert sanitize_filename("   ") == "export"


def test_response_chart_spec_normalizes_bar_line_pie() -> None:
    from core.export.response_render import response_chart_spec

    bar = response_chart_spec({"response_type": "chart_bar", "data": {"labels": ["A", "B"], "values": [1, 2]}})
    assert bar == {"kind": "bar", "labels": ["A", "B"], "values": [1.0, 2.0], "title": ""}

    pie = response_chart_spec({"response_type": "chart_pie", "data": {"labels": ["X"], "values": ["3"]}})
    assert pie["kind"] == "pie"
    assert pie["values"] == [3.0]

    assert response_chart_spec({"response_type": "data_table", "data": {}}) is None
    assert response_chart_spec({"response_type": "chart_bar", "data": {}}) is None


def test_response_scatter_spec_supports_points_and_xy() -> None:
    from core.export.response_render import response_scatter_spec

    via_points = response_scatter_spec({"response_type": "chart_scatter", "data": {"points": [[1, 2], [3, 4]]}})
    assert via_points["points"] == [(1.0, 2.0), (3.0, 4.0)]

    via_xy = response_scatter_spec({"response_type": "chart_scatter", "data": {"x": [1, 2], "y": [5, 6]}})
    assert via_xy["points"] == [(1.0, 5.0), (2.0, 6.0)]

    assert response_scatter_spec({"response_type": "chart_bar", "data": {}}) is None


def test_iter_sub_panels_reads_multi_panel_and_combined() -> None:
    from core.export.response_render import iter_sub_panels

    panels = iter_sub_panels(
        {"response_type": "multi_panel", "data": {"panels": [{"title": "A", "response_type": "kpi_card", "data": {}}]}}
    )
    assert panels is not None
    assert panels[0]["title"] == "A"

    assert iter_sub_panels({"response_type": "chart_bar", "data": {}}) is None


@pytest.mark.asyncio
async def test_pdf_export_renders_charts_tables_and_page_numbers() -> None:
    exporter = ExporterFactory.create("pdf")
    big_table_rows = [["Account", 100 + i] for i in range(250)]
    result = await exporter.export(
        payload=ExportPayload(
            title="Full Coverage Report",
            content_type="dashboard",
            data={
                "filters": {"region": "APAC"},
                "cards": [
                    {"title": "KPI", "response": {"response_type": "kpi_card", "data": {"label": "Users", "value": "42"}}},
                    {
                        "title": "Bar",
                        "response": {"response_type": "chart_bar", "data": {"labels": ["A", "B"], "values": [1, 2]}},
                    },
                    {
                        "title": "Line",
                        "response": {"response_type": "chart_line", "data": {"labels": ["A", "B"], "values": [1, 2]}},
                    },
                    {
                        "title": "Pie",
                        "response": {"response_type": "chart_pie", "data": {"labels": ["A", "B"], "values": [1, 2]}},
                    },
                    {
                        "title": "Scatter",
                        "response": {"response_type": "chart_scatter", "data": {"points": [[1, 2], [3, 4]]}},
                    },
                    {
                        "title": "Big table",
                        "response": {"response_type": "data_table", "data": {"columns": ["Account", "Value"], "rows": big_table_rows}},
                    },
                    {
                        "title": "Panels",
                        "response": {
                            "response_type": "multi_panel",
                            "data": {"panels": [{"title": "Sub", "response_type": "kpi_card", "data": {"label": "X", "value": 1}}]},
                        },
                    },
                ],
            },
        )
    )
    assert result.data.startswith(b"%PDF")
    # 250-row table plus charts should force pagination beyond the cover page.
    page_count = result.data.count(b"/Type /Page\n") + result.data.count(b"/Type/Page\n")
    assert page_count > 3
    assert len(result.data) > 2000


@pytest.mark.asyncio
async def test_pptx_export_builds_native_charts_and_tables() -> None:
    from pptx import Presentation

    exporter = ExporterFactory.create("pptx")
    result = await exporter.export(
        payload=ExportPayload(
            title="Deck With Charts",
            content_type="dashboard",
            data={
                "cards": [
                    {
                        "title": "Revenue",
                        "response": {"response_type": "chart_bar", "data": {"labels": ["A", "B"], "values": [1, 2]}},
                    },
                    {
                        "title": "Table",
                        "response": {
                            "response_type": "data_table",
                            "data": {"columns": ["Account", "MRR"], "rows": [["Acme", 100], ["Globex", 200]]},
                        },
                    },
                    {
                        "title": "No response",
                        "summary": "fallback summary text",
                    },
                ]
            },
        )
    )
    prs = Presentation(BytesIO(result.data))
    slides = list(prs.slides)
    assert len(slides) == 4  # title + 3 cards
    assert any(shape.has_chart for shape in slides[1].shapes)
    assert any(shape.has_table for shape in slides[2].shapes)
