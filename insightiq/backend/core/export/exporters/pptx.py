from __future__ import annotations

from datetime import UTC, datetime
from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from core.export.base import EXPORTERS, ExportPayload, ExportResult, IExporter
from core.export.response_render import (
    format_response_text,
    iter_sub_panels,
    response_chart_spec,
    response_scatter_spec,
    response_table_matrix,
    sanitize_filename,
)

BRAND_PRIMARY = RGBColor(0x25, 0x63, 0xEB)
BRAND_DARK = RGBColor(0x0F, 0x17, 0x2A)
BRAND_MUTED = RGBColor(0x64, 0x74, 0x8B)
BRAND_BORDER = RGBColor(0xE2, 0xE8, 0xF0)
BRAND_SURFACE = RGBColor(0xF8, 0xFA, 0xFC)
CHART_PALETTE = [
    RGBColor(0x25, 0x63, 0xEB),
    RGBColor(0x0E, 0xA5, 0xE9),
    RGBColor(0x22, 0xC5, 0x5E),
    RGBColor(0xF5, 0x9E, 0x0B),
    RGBColor(0xEF, 0x44, 0x44),
    RGBColor(0xA8, 0x55, 0xF7),
    RGBColor(0x14, 0xB8, 0xA6),
]

_CHART_TYPES = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
}

TITLE_ONLY_LAYOUT = 5
BLANK_LAYOUT = 6
MAX_TABLE_ROWS = 18


@EXPORTERS.register("pptx")
class PptxExporter(IExporter):
    async def export(self, *, payload: ExportPayload) -> ExportResult:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        _build_title_slide(prs, payload)

        cards = payload.data.get("cards", [])
        if not cards and payload.content_type == "conversation":
            cards = [
                {
                    "title": f"{m.get('role', 'user')}".title(),
                    "response": {"response_type": "explanation", "data": {"output": str(m.get("content", ""))}},
                }
                for m in payload.data.get("messages", [])
            ]

        for card in cards:
            response = card.get("response") or {}
            if not response and card.get("summary"):
                response = {"response_type": "explanation", "data": {"output": str(card.get("summary", ""))}}
            _build_card_slides(prs, title=str(card.get("title", "Card")), response=response)

        buffer = BytesIO()
        prs.save(buffer)
        stem = sanitize_filename(payload.title, default="insightiq_export")
        return ExportResult(
            filename=f"{stem}.pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            data=buffer.getvalue(),
        )


def _build_title_slide(prs: Presentation, payload: ExportPayload) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), prs.slide_width, Emu(120000))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BRAND_PRIMARY
    accent.line.fill.background()
    accent.shadow.inherit = False

    wordmark = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(6), Inches(0.6))
    wm_para = wordmark.text_frame.paragraphs[0]
    wm_run = wm_para.add_run()
    wm_run.text = "InsightIQ"
    wm_run.font.size = Pt(20)
    wm_run.font.bold = True
    wm_run.font.color.rgb = BRAND_PRIMARY

    title_ph = slide.shapes.title
    title_ph.text = payload.title[:140]
    for para in title_ph.text_frame.paragraphs:
        para.font.size = Pt(36)
        para.font.bold = True
        para.font.color.rgb = BRAND_DARK

    kind_label = "Conversation Export" if payload.content_type == "conversation" else "Dashboard Report"
    generated = datetime.now(UTC).strftime("%B %d, %Y at %H:%M UTC")
    active_filters = {k: v for k, v in (payload.data.get("filters") or {}).items() if v}
    subtitle_lines = [kind_label, f"Generated {generated}"]
    if active_filters:
        subtitle_lines.append("Filters: " + ", ".join(f"{k}={v}" for k, v in active_filters.items()))

    subtitle_ph = slide.placeholders[1]
    subtitle_ph.text_frame.clear()
    for i, line in enumerate(subtitle_lines):
        para = subtitle_ph.text_frame.paragraphs[0] if i == 0 else subtitle_ph.text_frame.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.size = Pt(14)
        run.font.color.rgb = BRAND_MUTED


def _new_content_slide(prs: Presentation, title: str) -> Any:
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY_LAYOUT])
    slide.shapes.title.text = title[:90]
    for para in slide.shapes.title.text_frame.paragraphs:
        para.font.size = Pt(24)
        para.font.bold = True
        para.font.color.rgb = BRAND_DARK
    return slide


def _build_card_slides(prs: Presentation, *, title: str, response: dict[str, Any]) -> None:
    sub_panels = iter_sub_panels(response)
    if sub_panels is not None:
        if not sub_panels:
            slide = _new_content_slide(prs, title)
            _add_text_box(slide, "No panels available.")
            return
        for panel in sub_panels:
            panel_title = str(panel.get("title") or panel.get("response_type", "Panel"))
            _build_card_slides(prs, title=f"{title} \u2014 {panel_title}", response=panel)
        return

    slide = _new_content_slide(prs, title)
    rtype = str(response.get("response_type", ""))
    data = response.get("data") or {}

    chart_spec = response_chart_spec(response)
    if chart_spec:
        _add_native_chart(slide, chart_spec)
        return

    scatter_spec = response_scatter_spec(response)
    if scatter_spec:
        _add_scatter_chart(slide, scatter_spec)
        return

    matrix = response_table_matrix(response)
    if matrix:
        _add_native_table(slide, matrix)
        return

    if rtype == "kpi_card":
        _add_kpi(slide, label=str(data.get("label", "Value")), value=str(data.get("value", "")))
        return

    text = format_response_text(response)
    _add_text_box(slide, text)


def _add_kpi(slide: Any, *, label: str, value: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    value_para = tf.paragraphs[0]
    value_run = value_para.add_run()
    value_run.text = value
    value_run.font.size = Pt(60)
    value_run.font.bold = True
    value_run.font.color.rgb = BRAND_PRIMARY
    value_para.alignment = PP_ALIGN.LEFT

    label_para = tf.add_paragraph()
    label_run = label_para.add_run()
    label_run.text = label
    label_run.font.size = Pt(18)
    label_run.font.color.rgb = BRAND_MUTED


def _add_text_box(slide: Any, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.3))
    tf = box.text_frame
    tf.word_wrap = True
    lines = (text or "").splitlines() or [""]
    for i, line in enumerate(lines[:60]):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.size = Pt(15)
        run.font.color.rgb = BRAND_DARK


def _add_native_table(slide: Any, matrix: list[list[str]]) -> None:
    header, *rows = matrix
    capped_rows = rows[:MAX_TABLE_ROWS]
    overflow_note = None
    if len(rows) > MAX_TABLE_ROWS:
        overflow_note = f"+{len(rows) - MAX_TABLE_ROWS} more rows (see full export for details)"

    n_rows = len(capped_rows) + 1
    n_cols = len(header)
    left, top, width, height = Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.42 * n_rows)
    height = min(height, Inches(5.4))
    graphic_frame = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = graphic_frame.table

    for c, col_name in enumerate(header):
        cell = table.cell(0, c)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_DARK
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for r, row in enumerate(capped_rows, start=1):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.text = str(row[c]) if c < len(row) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = BRAND_SURFACE if r % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(11)
                para.font.color.rgb = BRAND_DARK

    if overflow_note:
        note = slide.shapes.add_textbox(left, top + height + Inches(0.15), width, Inches(0.4))
        run = note.text_frame.paragraphs[0].add_run()
        run.text = overflow_note
        run.font.size = Pt(11)
        run.font.italic = True
        run.font.color.rgb = BRAND_MUTED


def _add_native_chart(slide: Any, spec: dict[str, Any]) -> None:
    kind = spec["kind"]
    labels = spec["labels"][:12]
    values = spec["values"][:12]

    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series(spec.get("title") or "Series 1", values)

    x, y, cx, cy = Inches(1.2), Inches(1.5), Inches(10.9), Inches(5.4)
    graphic_frame = slide.shapes.add_chart(_CHART_TYPES[kind], x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    chart.has_legend = kind == "pie"
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False

    plot = chart.plots[0]
    plot.has_data_labels = kind == "pie"
    series = plot.series[0]
    if kind == "pie":
        for i, point in enumerate(series.points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = CHART_PALETTE[i % len(CHART_PALETTE)]
    elif kind == "bar":
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = BRAND_PRIMARY
    else:
        series.format.line.color.rgb = BRAND_PRIMARY
        series.format.line.width = Pt(2.25)


def _add_scatter_chart(slide: Any, spec: dict[str, Any]) -> None:
    points = spec["points"][:200]
    chart_data = XyChartData()
    series = chart_data.add_series(spec.get("title") or "Series 1")
    for x_val, y_val in points:
        series.add_data_point(x_val, y_val)

    x, y, cx, cy = Inches(1.2), Inches(1.5), Inches(10.9), Inches(5.4)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.series[0].marker.format.fill.solid()
    plot.series[0].marker.format.fill.fore_color.rgb = BRAND_PRIMARY
